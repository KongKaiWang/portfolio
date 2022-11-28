"""
Jira Downtime Visualization

This program is designed to extract and visualize data from Jira Downtime Tracker

Authored by: Kong Wang
Date: 04/06/22

Notes:

Version Updates:
v1.0 (04/06/2022) - Initial release
v1.1 (04/11/2022) - Added additional visualizations stratified by process, month, statistical summary, etc.
v1.2 (04/12/2022) - Completed login check, added logging capabilities
v1.3 (04/13/2022) - Updated visuals with data labels, added send mail function
v1.4 (04/27/2022) - Added compliance assessment (%) graph. Updated extract Jira data function to remove duplicates.
                    Updated plot data function to select 'Suspended', 'Foil break', and 'Powered Off'
v1.5 (05/01/2022) - Added logic to update new password if required (still debugging). Fixed monthly data filter bug.
v1.6 (05/03/2022) - Added extract_compliance_data function. Updated extract_jira_data, extract_mes_data, and plot_data
                    functions to aggregate downtime by shift.
v1.7 (05/08/2022) - Updated function annotations
v1.8 (05/16/2022) - Updated criteria for function extract_mes_data to include 'Awaiting Current Process' for cathode
"""

import os
import glob
import time
import smtplib
import logging
import pandas as pd
import seaborn as sns
from email import encoders
from email.utils import formatdate
from email.mime.base import MIMEBase
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart
from pptx import Presentation
from pptx.util import Inches
from datetime import datetime, timedelta
from selenium import webdriver
from selenium.webdriver.common.by import By
from selenium.webdriver.chrome.service import Service
from selenium.webdriver.support.ui import WebDriverWait
from selenium.webdriver.support import expected_conditions as EC


def check_cfg(cfg: str) -> None:
    """
    Inputs: cfg(str) of file location for config file
    Function: Checks if config file exists - if not, requests user for information
    Returns: None
    """
    logging.info('Check login')
    if not os.path.exists(cfg):
        print('Config file does not exist.')
        # Request user info
        user_id = input('Please input your user ID')
        jira_pw = input('Please input your Jira password')
        mes_pw = input('Please input your MES password')
        # Save configuration information
        credentials = user_id + '\n' + jira_pw + '\n' + mes_pw
        with open(cfg, 'w') as file:
            logging.info(f'Writing configuration to {credentials}')
            file.write(credentials)


def download_jira_data(user_id: str,
                       password: str,
                       url: str) -> pd.DataFrame():
    """
    Inputs: user_id(str) for Jira login, password(str) for Jira login, url(str) for Jira site
    Function: downloads machine downtime data from Jira
    Returns: data_file(df) of downloaded Jira data
    """
    logging.info('Download Jira data')
    browser = webdriver.Chrome(service=Service(r'C:\Program Files (x86)\Google\Chrome\chromedriver.exe'))
    browser.get(url)
    # Login to Jira
    WebDriverWait(browser, 30).until(EC.element_to_be_clickable((By.XPATH, "//input[@id='login-form-username']")))
    browser.find_element(By.XPATH, "//input[@id='login-form-username']").send_keys(user_id)
    browser.find_element(By.XPATH, "//input[@id='login-form-password']").send_keys(password)
    browser.find_element(By.XPATH, "//input[@id='login-form-submit']").click()
    # Download Jira data
    WebDriverWait(browser, 30).until(EC.element_to_be_clickable(
        (By.XPATH, "//a[normalize-space()='View all issues and filters']")))
    browser.find_element(By.XPATH,
                         "//a[normalize-space()='View all issues and filters']").click()
    WebDriverWait(browser, 30).until(EC.element_to_be_clickable((By.XPATH, "//button[@id='AJS_DROPDOWN__76']")))
    browser.find_element(By.XPATH, "//button[@id='AJS_DROPDOWN__76']").click()
    WebDriverWait(browser, 30).until(EC.element_to_be_clickable((By.XPATH, "//a[@id='allCsvFields']")))
    browser.find_element(By.XPATH, "//a[@id='allCsvFields']").click()
    WebDriverWait(browser, 30).until(EC.element_to_be_clickable(
        (By.XPATH, "//button[@id='csv-export-dialog-export-button']")))
    browser.find_element(By.XPATH, "//button[@id='csv-export-dialog-export-button']").click()
    time.sleep(1)
    # Wait for download to complete
    files = glob.glob(rf'C:\Users\{user_id}\Downloads\*')
    files.sort(key=lambda x: os.path.getctime(x))
    while not files[-1].endswith('.csv') and not files[-1].startswith('Jira'):
        time.sleep(1)
        files = glob.glob(rf'C:\Users\{user_id}\Downloads\*')
        files.sort(key=lambda x: os.path.getctime(x))
    browser.quit()
    data_file = files[-1]
    return data_file


def extract_jira_data(file: str,
                      shift: pd.DataFrame()) -> pd.DataFrame():
    """
    Inputs: file(str) for data extraction, shift(dataframe) of shift schedule
    Function: extracts Jira data and saves to folder
    Returns: df(dataframe) of formatted Jira data
    """
    logging.info('Extract Jira data')
    # Read and format raw data file. Select important columns and compute downtime duration
    df = pd.read_csv(file, parse_dates=['Custom field (Start of Downtime)'])
    cols = ['Summary', 'Issue key', 'Issue id', 'Reporter', 'Creator', 'Created', 'Custom field (Action Taken)',
            'Custom field (Description of Issue)', 'Custom field (Downtime Reason)', 'Custom field (Group)',
            'Custom field (End of Downtime)', 'Custom field (Machine)', 'Custom field (Main Process)',
            'Custom field (Simple Root Cause (Electrode - Mixing))', 'Custom field (Start of Downtime)',
            'Custom field (Simple Root Cause (Electrode - Unplanned Downtime))',
            'Custom field (Simple Root Cause (Electrode - Planned Downtime))']
    df = df[cols]
    # Remove duplicate entries from submission updates
    df.sort_values(by='Created')
    df.dropna(subset=['Custom field (Start of Downtime)', 'Custom field (End of Downtime)'], inplace=True)
    df.drop_duplicates(subset='Custom field (Start of Downtime)', inplace=True)
    df.reset_index(inplace=True)
    dur_s = pd.to_datetime(df['Custom field (End of Downtime)']) - pd.to_datetime(
        df['Custom field (Start of Downtime)'])
    dur_s = dur_s.apply(lambda x: abs(x.total_seconds()/3600))
    df['Downtime Duration'] = dur_s
    # Combine downtime reason columns
    df['Downtime Reason'] = df['Custom field (Simple Root Cause (Electrode - Unplanned Downtime))'].fillna(
        df['Custom field (Simple Root Cause (Electrode - Mixing))']
        [df['Custom field (Simple Root Cause (Electrode - Mixing))'].notna()]).fillna(
        df['Custom field (Simple Root Cause (Electrode - Planned Downtime))']
        [df['Custom field (Simple Root Cause (Electrode - Planned Downtime))'].notna()])
    # df['Month'] = df['Custom field (Start of Downtime)'].apply(lambda x: pd.to_datetime(x).strftime('%B'))
    # Aggregate downtime by shift
    a, b, c, d = [[0] * len(df.index) for _ in range(4)]
    for i in range(len(df.index)):
        start = df['Custom field (Start of Downtime)'][i]
        end = df['Custom field (End of Downtime)'][i]
        shift_data = shift.loc[start:end][:-1].value_counts()
        # Populate downtime duration sum by shift
        for j in range(len(shift_data.index)):
            shift_label = shift_data.index[j][0]
            shift_duration = shift_data.values[j]
            if shift_label == 'A':
                a[i] = shift_duration
            elif shift_label == 'B':
                b[i] = shift_duration
            elif shift_label == 'C':
                c[i] = shift_duration
            elif shift_label == 'D':
                d[i] = shift_duration
    df['Jira A Shift'], df['Jira B Shift'], df['Jira C Shift'], df['Jira D Shift'] = a, b, c, d
    df.to_csv(rf'{DIRECTORY}\{MDY}\Jira_Data_{MDY}.csv', index=False)
    return df


def download_mes_data(**settings) -> pd.DataFrame():
    """
    Inputs: settings(dict) of query parameters
    Function: downloads machine status data from MES
    Returns: data_file(df) of downloaded MES data
    """
    logging.info('Download mes data')
    browser = webdriver.Chrome(service=Service(r'C:\Program Files (x86)\Google\Chrome\chromedriver.exe'))
    browser.get(settings['URL'])
    WebDriverWait(browser, 30).until(EC.element_to_be_clickable((By.XPATH, "//input[@id='submitButton']")))
    # Login to MES
    browser.find_element(By.XPATH, "//input[@id='txtUserId']").send_keys(settings['ID'])
    browser.find_element(By.XPATH, "//input[@id='txtPassword']").send_keys(settings['Password'])
    browser.find_element(By.XPATH, "//input[@id='submitButton']").click()
    # Update login if required
    if browser.find_element(By.XPATH, "//input[@id='ID_Dlg_txtPasswordNew']").is_displayed():
        new_pass = settings['Password'][:-2] + str(int(settings['Password'][-2]) + 1) + settings['Password'][-1]
        print(new_pass)
        with open(settings['Config'], 'w') as file:
            file.write(settings['ID'] + '\n' + settings['Password'] + '\n' + new_pass)
        WebDriverWait(browser, 30).until(EC.element_to_be_clickable(
            (By.XPATH, "// input[ @ id = 'ID_CM110_cmdUpdate']")))
        browser.find_element(By.XPATH, "//input[@id='ID_Dlg_txtPasswordNew']").send_keys(new_pass)
        browser.find_element(By.XPATH, "// input[ @ id = 'ID_Dlg_txtPasswordConfirm']").send_keys(new_pass)
        browser.find_element(By.XPATH, "// input[ @ id = 'ID_CM110_cmdUpdate']").click()
    # Download MES data
    WebDriverWait(browser, 30).until(EC.element_to_be_clickable(
        (By.XPATH, "//dt[normalize-space()='Production and Material Management']")))
    browser.find_element(By.XPATH,
                         "//dt[normalize-space()='Production and Material Management']").click()
    WebDriverWait(browser, 30).until(EC.element_to_be_clickable(
        (By.XPATH, "//input[@Value='RN200:Machine Status Report']")))
    browser.find_element(By.XPATH,
                         "//input[@value='RN200:Machine Status Report']").click()
    browser.switch_to.window(browser.window_handles[1])
    WebDriverWait(browser, 30).until(EC.element_to_be_clickable((By.XPATH, "//input[@id='ID_cmdCsv']")))
    browser.find_element(By.XPATH, "//select[@id='ID_cmbCndProcessGroupNo']").send_keys(settings['Group'])
    browser.find_element(By.XPATH, "//input[@id='ID_rdoSearchModeOper']").click()
    browser.find_element(By.XPATH, f"//option[@title='{settings['Process']}']").click()
    browser.find_element(By.XPATH, "//input[@id='ID_txtCndTermStart']").clear()
    browser.find_element(By.XPATH, "//input[@id='ID_txtCndTermStart']").send_keys(settings['Start'])
    browser.find_element(By.XPATH, "//input[@id='ID_txtCndTermEnd']").clear()
    browser.find_element(By.XPATH, "//input[@id='ID_txtCndTermEnd']").send_keys(settings['End'])
    browser.find_element(By.XPATH, "//input[@id='ID_cmdCsv']").click()
    WebDriverWait(browser, 30).until(EC.element_to_be_clickable(
        (By.XPATH, "//input[@id='ID_CSVOUTPUT_btnYes']")))
    browser.find_element(By.XPATH, "//input[@id='ID_CSVOUTPUT_chkOutputSearchCondition']").click()
    browser.find_element(By.XPATH, "//input[@id='ID_CSVOUTPUT_btnYes']").click()
    time.sleep(1)
    # Wait for download to complete
    files = glob.glob(rf"C:\Users\{settings['ID']}\Downloads\*")
    files.sort(key=lambda x: os.path.getctime(x))
    while not files[-1].endswith('.csv') and not files[-1].startswith('RN200'):
        time.sleep(1)
        files = glob.glob(rf"C:\Users\{settings['ID']}\Downloads\*")
        files.sort(key=lambda x: os.path.getctime(x))
    browser.quit()
    data_file = files[-1]
    return data_file


def extract_mes_data(anode: str,
                     cathode: str,
                     shift: pd.DataFrame()) -> pd.DataFrame():
    """
    Inputs: anode(str) file and cathode(str) file of MES runtime status, shift(dataframe) of shift schedule
    Function: extracts machine status and downtime data and saves to folder
    Returns: machine_downtime(dataframe), machine_status(dataframe) of formatted MES data
    """
    logging.info('Extract mes data')
    dfa = pd.read_csv(anode, parse_dates=['Start D/T', 'End D/T'])
    dfc = pd.read_csv(cathode, parse_dates=['Start D/T', 'End D/T'])
    # Aggregate downtime by shifts
    anode_downtime = dfa[dfa['Machine Status'].isin(['Suspended', 'Foil break', 'Powered Off'])].reset_index()
    cathode_downtime = dfc[dfc['Machine Status'].isin(['Suspended', 'Foil break', 'Powered Off',
                                                       'Awaiting Current Process'])].reset_index()
    for df in [anode_downtime, cathode_downtime]:
        a, b, c, d = [[0] * len(df.index) for _ in range(4)]
        for i in range(len(df.index)):
            start = df['Start D/T'][i]
            end = df['End D/T'][i]
            shift_data = shift.loc[start:end][:-1].value_counts()
            # Populate downtime duration sum by shift
            for j in range(len(shift_data.index)):
                shift_label = shift_data.index[j][0]
                shift_duration = shift_data.values[j]
                if shift_label == 'A':
                    a[i] = shift_duration
                elif shift_label == 'B':
                    b[i] = shift_duration
                elif shift_label == 'C':
                    c[i] = shift_duration
                elif shift_label == 'D':
                    d[i] = shift_duration
        df['MES A Shift'], df['MES B Shift'], df['MES C Shift'], df['MES D Shift'] = a, b, c, d
    machine_downtime = pd.concat([anode_downtime[anode_downtime.columns[-4:]].sum(axis=0),
                                  cathode_downtime[cathode_downtime.columns[-4:]].sum(axis=0)], axis=1)
    machine_downtime.columns = ['Anode', 'Cathode']
    machine_downtime = machine_downtime.T.apply(lambda x: x/60)
    machine_downtime.index.names = ['Polarity']
    machine_downtime.to_csv(rf'{DIRECTORY}\{MDY}\Machine_Downtime_{MDY}.csv')
    # Aggregate duration by machine status
    time_anode_hour = dfa.groupby('Machine Status').sum()['Elapsed Time(sec.)'].apply(lambda x: float(x)/3600)
    time_anode_percentage = dfa.groupby('Machine Status').sum()['Elapsed Time(sec.)'].apply(lambda x: float(x)/(36*1344))
    time_cathode_hour = dfc.groupby('Machine Status').sum()['Elapsed Time(sec.)'].apply(lambda x: float(x)/3600)
    time_cathode_percentage = dfc.groupby('Machine Status').sum()['Elapsed Time(sec.)'].apply(lambda x: float(x)/(36*1344))
    anode_status = pd.concat([time_anode_hour, time_anode_percentage], axis=1)
    cathode_status = pd.concat([time_cathode_hour, time_cathode_percentage], axis=1)
    anode_status.columns = ['Duration (h)', 'Duration (%)']
    cathode_status.columns = ['Duration (h)', 'Duration (%)']
    anode_status['Polarity'] = 'Anode'
    cathode_status['Polarity'] = 'Cathode'
    machine_status = pd.concat([anode_status, cathode_status])
    machine_status.reset_index(inplace=True)
    machine_status.to_csv(rf'{DIRECTORY}\{MDY}\Machine_Status_{MDY}.csv', index=False)
    return machine_status, machine_downtime


def extract_compliance_data(jira: pd.DataFrame(),
                            mes: pd.DataFrame()) -> pd.DataFrame():
    """
    Inputs: jira(dataframe) of Jira downtime data, mes(dataframe) of MES downtime data
    Function: extract compliance data from jira and mes dataframe
    Returns: df(dataframe) of compliance data
    """
    logging.info('Extract compliance data')
    # Format Jira data
    jira = jira[pd.to_datetime(jira['Custom field (Start of Downtime)']) >= START]
    jira_total = jira.groupby('Custom field (Main Process)').sum()[['Downtime Duration']]
    jira_shift = jira[['Custom field (Main Process)', 'Jira A Shift', 'Jira B Shift', 'Jira C Shift', 'Jira D Shift']]
    jira_shift_agg = jira_shift.groupby('Custom field (Main Process)').sum().apply(lambda x: x/60)
    # Aggregate MES & Jira totals
    mes_total = mes.sum(axis=1)
    df = pd.concat([jira_total, mes_total], axis=1)
    df.columns = ['Jira Downtime (h)', 'MES Downtime (h)']
    df['Jira %'] = df['Jira Downtime (h)'] * 100 / df['MES Downtime (h)']
    df['MES %'] = 100 - df['Jira %']
    # Aggregate downtime by shift
    df = pd.concat([df, mes], axis=1)
    df = pd.concat([df, jira_shift_agg], axis=1)
    # df['Jira A Shift %'] = df['Jira A Shift'] * 100 / df['MES A Shift']
    # df['Jira B Shift %'] = df['Jira B Shift'] * 100 / df['MES B Shift']
    # df['Jira C Shift %'] = df['Jira C Shift'] * 100 / df['MES C Shift']
    # df['Jira D Shift %'] = df['Jira D Shift'] * 100 / df['MES D Shift']
    # df['MES %'], df['MES A Shift %'], df['MES B Shift %'], df['MES C Shift %'], df['MES D Shift %'] = [100] * 5
    df.index.names = ['Polarity']
    df.to_csv(rf'{DIRECTORY}\{MDY}\Compliance_{MDY}.csv', index=False)
    return df


def plot_data(jira: pd.DataFrame(),
              mes_status: pd.DataFrame(),
              compliance: pd.DataFrame()) -> None:
    """
    Inputs: jira(dataframe) of Jira downtime data, mes_status(dataframe) of machine runtime, compliance(dataframe)
    Function: plots visualizations of downtime data and saves graphs to powerpoint
    Returns: None
    """
    logging.info('Plot data')
    # plot configurations
    palette = {'Anode': 'cyan', 'Cathode': 'purple'}
    colors = list(palette.values())
    jira_30d = jira[jira['Custom field (Start of Downtime)'] > (datetime.now() - timedelta(days=30))]
    anode_shift = pd.DataFrame(data={'MES': compliance.loc['Anode'][4:8].values,
                                     'Jira': compliance.loc['Anode'][-4:].values}, index=['A', 'B', 'C', 'D'])
    cathode_shift = pd.DataFrame(data={'MES': compliance.loc['Cathode'][4:8].values,
                                       'Jira': compliance.loc['Cathode'][-4:].values}, index=['A', 'B', 'C', 'D'])
    # downtime count by line (YTD)
    g1 = sns.catplot(data=jira.sort_values('Custom field (Machine)'), kind='count', x='Custom field (Machine)',
                     palette=palette, alpha=0.65, hue='Custom field (Main Process)', height=6)
    # downtime duration by line (YTD)
    g2 = sns.catplot(data=jira.sort_values('Custom field (Machine)'), kind='bar', x='Custom field (Machine)', ci=None,
                     y='Downtime Duration', palette=palette, alpha=0.65, hue='Custom field (Main Process)', height=6,
                     estimator=sum)
    # downtime count by line (30d rolling)
    g3 = sns.catplot(data=jira_30d.sort_values('Custom field (Machine)'),
                     kind='count', x='Custom field (Machine)', palette=palette, alpha=0.65,
                     hue='Custom field (Main Process)', height=6)
    # downtime duration by line (30d rolling)
    g4 = sns.catplot(data=jira_30d.sort_values('Custom field (Machine)'),
                     kind='bar', x='Custom field (Machine)', ci=None, y='Downtime Duration', palette=palette,
                     alpha=0.65, hue='Custom field (Main Process)', height=6, estimator=sum)
    # downtime count by reason (YTD anode)
    g5 = sns.catplot(data=jira[jira['Custom field (Main Process)'] == 'Anode'], kind='count', x='Downtime Reason',
                     palette=palette, alpha=0.65, height=6, hue='Custom field (Main Process)',
                     order=jira[jira['Custom field (Main Process)'] == 'Anode']['Downtime Reason'].value_counts().index)
    # downtime count by reason (YTD cathode)
    g6 = sns.catplot(data=jira[jira['Custom field (Main Process)'] == 'Cathode'], kind='count', x='Downtime Reason',
                     order=jira[jira['Custom field (Main Process)'] == 'Cathode']['Downtime Reason'].value_counts().index,
                     hue='Custom field (Main Process)', palette=palette, alpha=0.65, height=6)
    # downtime count by reason (30d rolling anode)
    g7 = sns.catplot(data=jira_30d[jira_30d['Custom field (Main Process)'] == 'Anode'], kind='count',
                     hue='Custom field (Main Process)', x='Downtime Reason', palette=palette, alpha=0.65, height=6,
                     order=jira_30d[jira_30d['Custom field (Main Process)'] == 'Anode'][
                         'Downtime Reason'].value_counts().index)
    # downtime count by reason (30d rolling cathode)
    g8 = sns.catplot(data=jira_30d[jira_30d['Custom field (Main Process)'] == 'Cathode'], kind='count',
                     hue='Custom field (Main Process)', x='Downtime Reason', palette=palette, alpha=0.65, height=6,
                     order=jira_30d[jira_30d['Custom field (Main Process)'] == 'Cathode'][
                         'Downtime Reason'].value_counts().index)
    # downtime duration by reason (YTD anode)
    g9 = sns.catplot(data=jira[jira['Custom field (Main Process)'] == 'Anode'].sort_values('Custom field (Machine)'),
                     y='Downtime Duration', alpha=0.65, height=6, ci=None, hue='Custom field (Main Process)',
                     kind='bar', x='Downtime Reason', palette=palette,
                     order=jira[jira['Custom field (Main Process)'] == 'Anode'].groupby(['Downtime Reason']).sum()[
                         'Downtime Duration'].sort_values(ascending=False).index, estimator=sum)
    # downtime duration by reason (YTD cathode)
    g10 = sns.catplot(data=jira[jira['Custom field (Main Process)'] == 'Cathode'].sort_values('Custom field (Machine)'),
                      y='Downtime Duration', alpha=0.65, height=6, ci=None, hue='Custom field (Main Process)',
                      kind='bar', x='Downtime Reason', palette=palette,
                      order=jira[jira['Custom field (Main Process)'] == 'Cathode'].groupby(['Downtime Reason']).sum()[
                          'Downtime Duration'].sort_values(ascending=False).index, estimator=sum)
    # downtime duration by reason (30d rolling anode)
    g11 = sns.catplot(data=jira_30d[jira_30d['Custom field (Main Process)'] == 'Anode'].sort_values('Custom field (Machine)'),
                      y='Downtime Duration', kind='bar', ci=None, hue='Custom field (Main Process)',
                      x='Downtime Reason', palette=palette, alpha=0.65, height=6, estimator=sum,
                      order=jira_30d[jira_30d['Custom field (Main Process)'] == 'Anode'].groupby(
                          ['Downtime Reason']).sum()['Downtime Duration'].sort_values(ascending=False).index)
    # downtime duration by reason (30d rolling cathode)
    g12 = sns.catplot(data=jira_30d[jira_30d['Custom field (Main Process)'] == 'Cathode'].sort_values('Custom field (Machine)'),
                      x='Downtime Reason', palette=palette, alpha=0.65, height=6, estimator=sum,
                      kind='bar', ci=None, y='Downtime Duration', hue='Custom field (Main Process)',
                      order=jira_30d[jira_30d['Custom field (Main Process)'] == 'Cathode'].groupby(
                          ['Downtime Reason']).sum()['Downtime Duration'].sort_values(ascending=False).index)
    # run status distribution (7d rolling)
    g13 = sns.catplot(data=mes_status, kind='bar', x='Machine Status', y='Duration (%)', hue='Polarity', height=6,
                      ci=None, palette=palette, alpha=0.65)
    # total compliance hours (7d rolling)
    g14_ax = compliance[['Jira Downtime (h)', 'MES Downtime (h)']].plot(
        kind='bar', title='Compliance Hrs (7d rolling)', ylabel='Compliance (Hrs)', stacked=True, color=colors,
        alpha=0.45)
    g14 = g14_ax.figure
    # total compliance percentage (7d rolling)
    g15_ax = compliance[['Jira %', 'MES %']].plot(kind='bar', title='Compliance % 7d rolling', alpha=0.45,
                                                  ylabel='Compliance (%)', stacked=True, color=colors)
    g15 = g15_ax.figure
    # anode shift compliance (7d rolling)
    g16_ax = anode_shift.plot(kind='bar', title='Shift Compliance Hrs 7d rolling (Anode)', alpha=0.45,
                              color=list(reversed(colors)), xlabel='Shift', ylabel='Downtime (h)')
    g16 = g16_ax.figure
    # cathode shift compliance (7d rolling)
    g17_ax = cathode_shift.plot(kind='bar', title='Shift Compliance Hrs 7d rolling (Cathode)', alpha=0.45,
                                color=list(reversed(colors)), xlabel='Shift', ylabel='Downtime (h)')
    g17 = g17_ax.figure
    # plot formatting
    g1.set_axis_labels('Line', 'Instances (count)')
    g1.fig.suptitle('Downtime Instances YTD')
    g2.set_axis_labels('Line', 'Duration (h)')
    g2.fig.suptitle('Downtime Duration YTD')
    g3.set_axis_labels('Line', 'Instances (count)')
    g3.fig.suptitle(f'Downtime Instances 30d Rolling')
    g4.set_axis_labels('Line', 'Duration (h)')
    g4.fig.suptitle(f'Downtime Duration 30d Rolling')
    g5.set_axis_labels('Downtime Reason', 'Instances (count)')
    g5.fig.suptitle('Downtime Instance Reasons YTD (Anode)')
    g5.set_xticklabels(rotation=90)
    g6.set_axis_labels('Downtime Reason', 'Instances (count)')
    g6.fig.suptitle('Downtime Instance Reasons YTD (Cathode)')
    g6.set_xticklabels(rotation=90)
    g7.set_axis_labels('Downtime Reason', 'Instances (count)')
    g7.fig.suptitle(f'Downtime Instance Reasons 30d Rolling (Anode)')
    g7.set_xticklabels(rotation=90)
    g8.set_axis_labels('Downtime Reason', 'Instance (h)')
    g8.fig.suptitle(f'Downtime Instance Reasons 30d Rolling (Cathode)')
    g8.set_xticklabels(rotation=90)
    g9.set_axis_labels('Downtime Reason', 'Duration (h)')
    g9.fig.suptitle('Downtime Duration Reasons YTD (Anode)')
    g9.set_xticklabels(rotation=90)
    g10.set_axis_labels('Downtime Reason', 'Duration (h)')
    g10.fig.suptitle('Downtime Duration Reasons YTD (Cathode)')
    g10.set_xticklabels(rotation=90)
    g11.set_axis_labels('Downtime Reason', 'Duration (h)')
    g11.fig.suptitle(f'Downtime Duration Reasons 30d Rolling (Anode)')
    g11.set_xticklabels(rotation=90)
    g12.set_axis_labels('Downtime Reason', 'Duration (h)')
    g12.fig.suptitle(f'Downtime Duration Reasons 30d Rolling (Cathode)')
    g12.set_xticklabels(rotation=90)
    g13.set_axis_labels('Machine Status', 'Time Distribution (%)')
    g13.fig.suptitle('Machine Run Status (7d rolling)')
    g13.set_xticklabels(rotation=90)
    g14_ax.bar_label(g14_ax.containers[0], fmt='%.0f', size=8)
    g14_ax.legend(loc='upper left', bbox_to_anchor=(1.04, 1))
    g14.tight_layout()
    g15_ax.bar_label(g15_ax.containers[0], fmt='%.0f', size=8)
    g15_ax.legend(loc='upper left', bbox_to_anchor=(1.04, 1))
    g15.tight_layout()
    g16_ax.bar_label(g16_ax.containers[0], fmt='%.0f', size=8)
    g16_ax.bar_label(g16_ax.containers[1], fmt='%.0f', size=8)
    g16_ax.legend(loc='upper left', bbox_to_anchor=(1.04, 1))
    g16.tight_layout()
    g17_ax.bar_label(g17_ax.containers[0], fmt='%.0f', size=8)
    g17_ax.bar_label(g17_ax.containers[1], fmt='%.0f', size=8)
    g17_ax.legend(loc='upper left', bbox_to_anchor=(1.04, 1))
    g17.tight_layout()
    g_list = [g1, g2, g3, g4, g5, g6, g7, g8, g9, g10, g11, g12, g13, g14, g15, g16, g17]
    for g in g_list[:-4]:
        g.tight_layout()
        ax = g.facet_axis(0, 0)
        if g not in [g5, g6, g9, g10]:
            for c in ax.containers:
                labels = [f'{(v.get_height()):.0f}' for v in c]
                ax.bar_label(c, labels=labels, label_type='edge', size=8)
    # save plots
    png_list = [
        rf'{DIRECTORY}\{MDY}\dt_count_ytd_{MDY}.png',
        rf'{DIRECTORY}\{MDY}\dt_duration_ytd_{MDY}.png',
        rf'{DIRECTORY}\{MDY}\dt_count_month_{MDY}.png',
        rf'{DIRECTORY}\{MDY}\dt_duration_month_{MDY}.png',
        rf'{DIRECTORY}\{MDY}\dt_count_reason_anode_ytd_{MDY}.png',
        rf'{DIRECTORY}\{MDY}\dt_count_reason_cathode_ytd_{MDY}.png',
        rf'{DIRECTORY}\{MDY}\dt_count_reason_anode_month_{MDY}.png',
        rf'{DIRECTORY}\{MDY}\dt_count_reason_cathode_month_{MDY}.png',
        rf'{DIRECTORY}\{MDY}\dt_duration_reason_anode_ytd_{MDY}.png',
        rf'{DIRECTORY}\{MDY}\dt_duration_reason_cathode_ytd_{MDY}.png',
        rf'{DIRECTORY}\{MDY}\dt_duration_reason_anode_month_{MDY}.png',
        rf'{DIRECTORY}\{MDY}\dt_duration_reason_cathode_month_{MDY}.png',
        rf'{DIRECTORY}\{MDY}\runtime_distribution_anode_{MDY}.png',
        rf'{DIRECTORY}\{MDY}\compliance_assessment_{MDY}.png',
        rf'{DIRECTORY}\{MDY}\compliance_percentage_{MDY}.png',
        rf'{DIRECTORY}\{MDY}\shift_compliance_anode_{MDY}.png',
        rf'{DIRECTORY}\{MDY}\shift_compliance_cathode_{MDY}.png'
    ]
    for g, png in zip(g_list, png_list):
        g.savefig(png)
    # create powerpoint
    pres = Presentation()
    for png in png_list:
        layout = pres.slide_layouts[0]
        slide = pres.slides.add_slide(layout)
        slide.shapes.add_picture(png, left=Inches(0.0), top=Inches(0), width=Inches(10), height=Inches(7.5))
    pres.save(rf'{DIRECTORY}\{MDY}\Downtime_Visuals_{MDY}.pptx')


def send_mail(send_from: str,
              send_to: list,
              subject: str,
              text: str,
              files: list,
              server: str,
              port: int) -> None:
    """
    Inputs: send_from(str), send_to(list), subject(str), text(str), files(list), server(str), port(int) - config info
    Function: Sends mail with downtime report attachment
    Returns: None
    """
    logging.info('Send mail')
    msg = MIMEMultipart()
    msg['From'] = send_from
    msg['To'] = ", ".join(send_to)
    msg['Date'] = formatdate(localtime=True)
    msg['Subject'] = subject
    msg.attach(MIMEText(text))
    for file in files:
        count = 0
        if not os.path.exists(file):
            time.sleep(1)
            count += 1
            if count >= 30:
                logging.info('Mail attachment files do not exist.')
                exit()
        part = MIMEBase('application', "octet-stream")
        part.set_payload(open(file, 'rb').read())
        encoders.encode_base64(part)
        part.add_header('Content-Disposition', 'attachment; filename={}'.format(os.path.split(file)[1]))
        msg.attach(part)
    smtp = smtplib.SMTP(server, port)
    smtp.sendmail(send_from, send_to, msg.as_string())
    smtp.quit()


def create_shift_schedule() -> None:
    """
    Inputs: None
    Function: Set up shift schedule reference .csv file
    Returns: None
    """
    minute_data = [datetime(2022, 1, 1) + timedelta(minutes=_) for _ in range(525600)]
    a_shift = [0, 4, 5, 8, 9, 10, 13]
    b_shift = [1, 2, 3, 4, 9, 10, 11]
    c_shift = [1, 2, 3, 6, 7, 11, 12]
    d_shift = [0, 5, 6, 7, 8, 12, 13]
    day_hrs = [_ for _ in range(7, 19)]
    night_hrs = [_ for _ in range(0, 7)] + [_ for _ in range(19, 24)]
    shift = ['B']
    for minute in minute_data[1:]:
        day_of_year = minute.timetuple().tm_yday
        hour_of_day = minute.hour
        if hour_of_day in day_hrs:
            if day_of_year % 14 in a_shift:
                shift.append('A')
            elif day_of_year % 14 in c_shift:
                shift.append('C')
        elif hour_of_day in night_hrs:
            # day change bug fix
            if shift[-1] in ['B', 'D']:
                shift.append(shift[-1])
            else:
                if day_of_year % 14 in b_shift:
                    shift.append('B')
                elif day_of_year % 14 in d_shift:
                    shift.append('D')
    shift_schedule = pd.DataFrame(data={'Time': minute_data, 'Shift': shift})
    shift_schedule.set_index('Time', inplace=True)
    shift_schedule.to_csv(rf'{DIRECTORY}\Shift_Schedule_test.csv')


if __name__ == '__main__':
    # Initialize parameters
    pd.set_option('display.width', 200, 'display.max_columns', 20)
    MDY = datetime.strftime(datetime.now(), '%m%d%Y')
    DIRECTORY = r'H:\BatteryEngineering\Electrode\Kong_Wang\Jira_Downtime'
    JIRA_URL = r'http://penajira/projects/ECD/issues/ECD-562?filter=allopenissues'
    MES_URL = r'http://gf-mes.america.gds.panasonic.com/GCAMES/'
    CFG = r'C:\Users\KW38770\Documents\WangK\Scripts\pena_credentials.txt'
    CHROMEDRIVER = r'C:\Program Files (x86)\Google\Chrome\chromedriver.exe'
    USER_ID, JIRA_PASSWORD, MES_PASSWORD = open(CFG).read().splitlines()
    START = datetime.strftime(datetime.now() - timedelta(days=7), '%m/%d/%Y')
    END = datetime.strftime(datetime.today(), '%m/%d/%Y')
    CATHODE = {'Group': 'Cathode', 'Process': 'EC-Coat', 'Start': START, 'End': END, 'ID': USER_ID,
               'Password': MES_PASSWORD, 'URL': MES_URL, 'Config': CFG}
    ANODE = {'Group': 'Anode', 'Process': 'EA-Coat', 'Start': START, 'End': END, 'ID': USER_ID,
             'Password': MES_PASSWORD, 'URL': MES_URL, 'Config': CFG}
    SHIFT = pd.read_csv(rf'{DIRECTORY}/Shift_Schedule.csv', index_col='Time', parse_dates=['Time'])
    # Mail settings
    SEND_FROM = 'Battery.Engineering@us.panasonic.com'
    SEND_TO = ['shane.thomas@us.panasonic.com', 'kevin.maningas@us.panasonic.com', 'kong.wang@us.panasonic.com']
    SUBJECT = 'Jira Downtime report'
    TEXT = 'Script ran successfully.'
    FILES = [rf'{DIRECTORY}\{MDY}\Downtime_Visuals_{MDY}.pptx']
    SERVER = 'mailrelay1.us.panasonic.com'
    PORT = 25

    # # 5/3/22 Unit Testing
    # jira_downtime = pd.read_csv(rf'{DIRECTORY}\{MDY}\Jira_Data_{MDY}.csv', parse_dates=['Custom field (Start of Downtime)'])
    # mes_machine_status = pd.read_csv(rf'{DIRECTORY}\{MDY}\Machine_Status_{MDY}.csv')
    # mes_machine_downtime = pd.read_csv(rf'{DIRECTORY}\{MDY}\Machine_Downtime_{MDY}.csv', index_col=['Polarity'])
    # compliance = extract_compliance_data(jira_downtime, mes_machine_downtime)
    # plot_data(jira_downtime, mes_machine_status, compliance)
    # # 5/2/22 Unit Testing
    # anode_mes_file = r'C:\Users\KW38770\Downloads\RN200_05022022132951313.csv'
    # cathode_mes_file = r'C:\Users\KW38770\Downloads\RN200_05022022133001032.csv'
    # mes_machine_status, mes_machine_downtime = extract_mes_data(anode_mes_file, cathode_mes_file, SHIFT)
    # jira_downtime = pd.read_csv(rf'{DIRECTORY}\{MDY}\Jira_Data_{MDY}.csv')
    # compliance = extract_compliance_data(jira_downtime, mes_machine_downtime)
    # plot_data(jira_downtime, mes_machine_status, compliance)  # plot data
    # 5/1/22 Unit Testing
    # plot_data(pd.read_csv(rf'{DIRECTORY}\{MDY}\Jira_Data_{MDY}.csv', parse_dates=['Custom field (Start of Downtime)']),
    # pd.read_csv(rf'{DIRECTORY}\{MDY}\MES_Reason_{MDY}.csv'))  # plot data
    # Initialize logger settings
    logging.basicConfig(filename=rf'{DIRECTORY}/Logs/{MDY}.log', level=logging.INFO,
                        format='%(asctime)s %(levelname)s %(name)s %(message)s')
    logger = logging.getLogger(__name__)

    # Main code
    try:
        logging.info('Initiating Jira Downtime Program - ' + str(datetime.now().time()))
        check_cfg(CFG)
        if not os.path.exists(CHROMEDRIVER):
            logging.info('Selenium chrome driver does not exist.')
            exit()
        if not os.path.exists(f'{DIRECTORY}\{MDY}'):
            os.mkdir(f'{DIRECTORY}\{MDY}')
        jira_data = download_jira_data(USER_ID, JIRA_PASSWORD, JIRA_URL)
        jira_downtime = extract_jira_data(jira_data, SHIFT)
        anode_mes_file = download_mes_data(**ANODE)
        cathode_mes_file = download_mes_data(**CATHODE)
        mes_machine_status, mes_machine_downtime = extract_mes_data(anode_mes_file, cathode_mes_file, SHIFT)
        compliance_data = extract_compliance_data(jira_downtime, mes_machine_downtime)
        plot_data(jira_downtime, mes_machine_status, compliance_data)
        send_mail(SEND_FROM, SEND_TO, SUBJECT, TEXT, FILES, SERVER, PORT)
        logging.info('Completed Jira Downtime Program - ' + str(datetime.now().time()))
    except Exception as err:
        logger.error(err)
